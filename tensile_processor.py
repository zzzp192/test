import docx
import openpyxl 
import statistics
import os
import re
import traceback 
from pptx import Presentation
from ppt_utils import (
    to_float, delete_table_column, delete_table_row, insert_table_row as insert_row,
    duplicate_slide, format_cell, clean_cell_merge_info, THEME_COLOR, RED_COLOR
)

# ================= 1. 数据提取模块 =================

def extract_from_docx(docx_path):
    doc = docx.Document(docx_path)
    project_id = "未知项目"
    try:
        if doc.sections[0].header.paragraphs:
            txt = doc.sections[0].header.paragraphs[0].text.strip()
            if txt: project_id = txt
    except: pass
    
    if project_id == "未知项目" and doc.paragraphs:
        txt = doc.paragraphs[0].text.split("：")[0].strip()
        project_id = txt

    extracted_groups = {} 
    if not doc.tables: return project_id, {}
    table = doc.tables[0]
    
    for row in table.rows[2:]:
        cells = row.cells
        if len(cells) < 13: continue 
        full_id = cells[1].text.strip()
        if not full_id: continue

        clean_id = re.sub(r'[\(（].*?[\)）]', '', full_id).strip()
        has_note = (clean_id != full_id)

        if "-" in clean_id:
            parts = clean_id.rsplit("-", 1)
            group_name = parts[0]
            sample_num = parts[1]
        else:
            group_name = clean_id
            sample_num = "1"

        item = {
            "id_num": sample_num,
            "thick": cells[3].text.strip(), # Word保持不变，如果Word也要改B列需确认位置
            "Rp": to_float(cells[8].text),
            "Rm": to_float(cells[9].text),
            "Ag": to_float(cells[10].text),
            "A":  to_float(cells[12].text),
            "has_note": has_note
        }

        if group_name not in extracted_groups:
            extracted_groups[group_name] = []
        extracted_groups[group_name].append(item)

    return project_id, extracted_groups

def extract_from_excel(xlsx_path):
    project_id = os.path.splitext(os.path.basename(xlsx_path))[0]
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    
    if "Sheet1" in wb.sheetnames:
        sheet = wb["Sheet1"]
    else:
        sheet = wb.worksheets[0]
    
    extracted_groups = {}
    
    for row in sheet.iter_rows(min_row=1, values_only=True):
        if not row or row[0] is None: continue
        
        raw_id = str(row[0]).strip()
        clean_id = re.sub(r'[\(（].*?[\)）]', '', raw_id).strip()
        has_note = (clean_id != raw_id)

        if "-" not in clean_id or not re.search(r'\d$', clean_id):
            continue

        parts = clean_id.rsplit("-", 1)
        group_name = parts[0]
        sample_num = parts[1]
        
        try:
            # --- 修改点 1: 抓取 B 列 (Index 1) 作为宽度/厚度 ---
            thick_val = row[1] if len(row) > 1 else "" 
            # -----------------------------------------------
            
            rp_val = row[6] if len(row) > 6 else 0
            rm_val = row[7] if len(row) > 7 else 0
            ag_val = row[8] if len(row) > 8 else 0
            a_val  = row[10] if len(row) > 10 else 0
            
            item = {
                "id_num": sample_num,
                "thick": str(thick_val) if thick_val is not None else "",
                "Rp": to_float(rp_val),
                "Rm": to_float(rm_val),
                "Ag": to_float(ag_val),
                "A":  to_float(a_val),
                "has_note": has_note
            }
            
            if group_name not in extracted_groups:
                extracted_groups[group_name] = []
            extracted_groups[group_name].append(item)
            
        except Exception as e:
            print(f"Skipping row {raw_id}: {e}")
            continue

    wb.close()
    return project_id, extracted_groups

def calculate_stats(group_data):
    if not group_data: return {}
    stats = {}
    for key in ['Rp', 'Rm']:
        vals = [d[key] for d in group_data]
        if len(vals) == 0:
             stats[key] = "/"
             continue
        m = statistics.mean(vals)
        s = statistics.stdev(vals) if len(vals) > 1 else 0.0
        stats[key] = f"{m:.0f}±{s:.0f}"
    for key in ['Ag', 'A']:
        vals = [d[key] for d in group_data]
        if len(vals) == 0:
             stats[key] = "/"
             continue
        m = statistics.mean(vals)
        s = statistics.stdev(vals) if len(vals) > 1 else 0.0
        stats[key] = f"{m:.1f}±{s:.1f}"
    return stats

# ================= 2. 主流程控制器 =================

def generate_report(data_path, pptx_template_path, output_path, include_ag=True):
    try:
        # 1. 提取
        ext = os.path.splitext(data_path)[1].lower()
        if ext == ".docx":
            project_id, groups = extract_from_docx(data_path)
        elif ext == ".xlsx" or ext == ".xls":
            project_id, groups = extract_from_excel(data_path)
        else:
            return "错误：不支持的文件格式"

        if not groups: return "错误：未提取到数据"

        # 2. PPT 处理
        prs = Presentation(pptx_template_path)
        group_names = list(groups.keys())
        
        tables = [s.table for s in prs.slides[0].shapes if s.has_table]
        if not tables: return "错误：PPT第1页未找到表格"
        template_table = tables[0]
        
        # 计算每页容量
        template_capacity = 0
        for row in template_table.rows:
            txt = "".join([c.text_frame.text for c in row.cells]).lower()
            if "平均" in txt or "average" in txt:
                template_capacity += 1
        if template_capacity == 0: template_capacity = 1
        
        GROUPS_PER_SLIDE = template_capacity
        total_slides = (len(group_names) + GROUPS_PER_SLIDE - 1) // GROUPS_PER_SLIDE
        
        # 复制页面
        for i in range(total_slides - 1):
            duplicate_slide(prs, 0)
            
        # 填充
        for slide_idx in range(total_slides):
            slide = prs.slides[slide_idx]
            # 填写项目号
            for shape in slide.shapes:
                if shape.has_text_frame and "项目号" in shape.text:
                    try: shape.text_frame.text = shape.text.replace("项目号", project_id)
                    except: pass

            slide_tables = [s.table for s in slide.shapes if s.has_table]
            if not slide_tables: continue
            main_table = slide_tables[0]

            # --- 修改点 2: 如果不包含 Ag，删除第 6 列 (Index 5) ---
            # 列索引: 0=Group, 1=ID, 2=Thick, 3=Rp, 4=Rm, 5=Ag, 6=A
            if not include_ag:
                # 只有当表格确实有这么多列时才删除
                if len(main_table.columns) >= 6:
                    delete_table_column(main_table, 5) 
            # --------------------------------------------------

            batch_groups = group_names[slide_idx*GROUPS_PER_SLIDE : (slide_idx+1)*GROUPS_PER_SLIDE]
            
            # 查找统计行位置 (需要重新查找，因为行数是动态的)
            avg_row_indices = []
            for r, row in enumerate(main_table.rows):
                txt = "".join([c.text_frame.text for c in row.cells]).lower()
                if "平均" in txt or "average" in txt:
                    avg_row_indices.append(r)
            
            # 倒序处理每一块
            for i in range(len(avg_row_indices) - 1, -1, -1):
                if i >= GROUPS_PER_SLIDE: continue
                
                # 重新定位统计行 (因为上面插入/删除行会导致索引变化)
                current_avg_rows = []
                for r, row in enumerate(main_table.rows):
                    txt = "".join([c.text_frame.text for c in row.cells]).lower()
                    if "平均" in txt or "average" in txt:
                        current_avg_rows.append(r)
                
                if i >= len(current_avg_rows): break
                stat_row_idx = current_avg_rows[i]
                
                if i == 0: start_row_idx = 1
                else: start_row_idx = current_avg_rows[i-1] + 1
                
                current_slots = stat_row_idx - start_row_idx
                
                if i < len(batch_groups):
                    g_name = batch_groups[i]
                    data = groups[g_name]
                    stats = calculate_stats(data)
                    n_samples = len(data)
                    
                    # 动态增减行
                    if n_samples > current_slots:
                        for _ in range(n_samples - current_slots):
                            insert_row(main_table, stat_row_idx, start_row_idx)
                            stat_row_idx += 1 
                    elif n_samples < current_slots:
                        for _ in range(current_slots - n_samples):
                            delete_table_row(main_table, stat_row_idx - 1)
                            stat_row_idx -= 1
                    
                    # 清理合并
                    for k in range(n_samples):
                        target_row = main_table.rows[start_row_idx + k]
                        clean_cell_merge_info(target_row.cells[0]) 

                    # 填数据
                    for k in range(n_samples):
                        row = main_table.rows[start_row_idx + k]
                        item = data[k]
                        font_color = RED_COLOR if item.get('has_note') else None
                        
                        format_cell(row.cells[0], g_name if k==0 else "")
                        format_cell(row.cells[1], str(item['id_num']), color_rgb=font_color)
                        format_cell(row.cells[2], str(item['thick']))
                        
                        # --- 修改点 3: 动态填充 ---
                        # 确保单元格索引不越界
                        c_idx = 3
                        if c_idx < len(row.cells): format_cell(row.cells[c_idx], str(item['Rp']))
                        c_idx += 1
                        if c_idx < len(row.cells): format_cell(row.cells[c_idx], str(item['Rm']))
                        c_idx += 1
                        
                        if include_ag:
                            if c_idx < len(row.cells): format_cell(row.cells[c_idx], str(item['Ag']))
                            c_idx += 1
                        
                        if c_idx < len(row.cells): format_cell(row.cells[c_idx], str(item['A']))
                        # -------------------------
                    
                    # 合并组名
                    if n_samples > 1:
                        try:
                            cell_top = main_table.cell(start_row_idx, 0)
                            cell_bottom = main_table.cell(start_row_idx + n_samples - 1, 0)
                            cell_top.merge(cell_bottom)
                            format_cell(cell_top, g_name)
                        except:
                            format_cell(main_table.cell(start_row_idx, 0), g_name)

                    # 填统计行
                    stat_row = main_table.rows[stat_row_idx]
                    format_cell(stat_row.cells[0], "平均值±标准差", is_bold=True, color_rgb=THEME_COLOR)
                    
                    # --- 修改点 4: 统计行动态填充 ---
                    s_idx = 3
                    if s_idx < len(stat_row.cells): format_cell(stat_row.cells[s_idx], stats['Rp'], is_bold=True, color_rgb=THEME_COLOR)
                    s_idx += 1
                    if s_idx < len(stat_row.cells): format_cell(stat_row.cells[s_idx], stats['Rm'], is_bold=True, color_rgb=THEME_COLOR)
                    s_idx += 1
                    
                    if include_ag:
                        if s_idx < len(stat_row.cells): format_cell(stat_row.cells[s_idx], stats['Ag'], is_bold=True, color_rgb=THEME_COLOR)
                        s_idx += 1
                        
                    if s_idx < len(stat_row.cells): format_cell(stat_row.cells[s_idx], stats['A'], is_bold=True, color_rgb=THEME_COLOR)
                    # ------------------------------

        prs.save(output_path)
        return f"成功！共处理 {len(groups)} 组数据。\n保存至: {output_path}"

    except Exception as e:
        return f"运行出错: {str(e)}\n\n详细位置:\n{traceback.format_exc()}"
