#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
育材堂报告助手 V3.7 - PPT操作工具模块

软件名称：育材堂报告助手
版本号：V3.7
开发单位：育材堂
开发完成日期：2024年

模块功能：
    提供PowerPoint文件操作的共享工具函数。

主要功能：
    - 表格行列的增删操作
    - 幻灯片复制功能
    - 单元格格式化和文本替换
    - 单元格合并信息清理

Copyright (c) 2024 育材堂. All rights reserved.
"""

import copy
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# 主题颜色常量
THEME_COLOR = RGBColor(25, 137, 141)
RED_COLOR = RGBColor(255, 0, 0)

def to_float(val):
    """安全转换为浮点数"""
    try:
        return float(val) if val is not None else 0.0
    except:
        return 0.0

def delete_table_column(table, col_idx):
    """
    删除PPT表格中的某一列（增强版：支持合并单元格）
    """
    tbl = table._tbl
    
    # 1. 删除物理列定义 (gridCol)
    tblGrid = tbl.tblGrid
    gridCols = tblGrid.findall(qn("a:gridCol"))
    if col_idx < len(gridCols):
        tblGrid.remove(gridCols[col_idx])
    
    # 2. 遍历每一行，寻找对应的单元格进行删除或缩减跨度
    for tr in tbl.findall(qn("a:tr")):
        tcs = tr.findall(qn("a:tc"))
        
        visual_col_idx = 0  # 当前遍历到的视觉列索引
        target_tc = None    # 目标单元格
        is_in_span = False  # 是否处于合并单元格内部
        
        # 寻找对应 col_idx 的单元格
        for tc in tcs:
            # 获取该单元格跨越的列数 (默认为1)
            grid_span = 1
            if tc.tcPr is not None:
                gs = tc.tcPr.find(qn("a:gridSpan"))
                if gs is not None:
                    grid_span = int(gs.get("val"))
            
            # 判断目标列是否在这个单元格范围内
            # 例如：单元格跨越 0-2 列，我们要删第 1 列，则命中
            if visual_col_idx <= col_idx < visual_col_idx + grid_span:
                target_tc = tc
                if grid_span > 1:
                    is_in_span = True
                break
            
            visual_col_idx += grid_span
            
        # 执行删除或修改操作
        if target_tc is not None:
            if is_in_span:
                # 如果目标列在合并单元格内，不删除单元格，而是将跨度减 1
                gs = target_tc.tcPr.find(qn("a:gridSpan"))
                current_span = int(gs.get("val"))
                if current_span > 1:
                    gs.set("val", str(current_span - 1))
            else:
                # 普通单元格，直接删除
                tr.remove(target_tc)

def delete_table_row(table, row_idx):
    """删除表格行"""
    if row_idx < 0 or row_idx >= len(table.rows):
        return
    tr = table.rows[row_idx]._tr
    tr.getparent().remove(tr)

def insert_table_row(table, target_idx, source_idx):
    """
    在指定位置插入新行（修复版）
    修复：复制时仅清除垂直合并(vMerge)，必须保留水平合并(gridSpan)以维持表格结构完整性。
    """
    tbl = table._tbl
    source_tr = table.rows[source_idx]._tr
    new_tr = copy.deepcopy(source_tr)
    
    for tc in new_tr.tc_lst:
        if tc.tcPr is not None:
            # === 关键修改 ===
            # 只移除 vMerge (垂直合并状态)，让新行作为新数据的开始
            # 绝对不要移除 gridSpan！否则会导致行宽与表格定义不匹配，造成文件损坏
            v_merge = tc.tcPr.find(qn("a:vMerge"))
            if v_merge is not None:
                tc.tcPr.remove(v_merge)
            
            # 清理文本内容，防止复制旧数据
            if hasattr(tc, 'txBody'):
                # 简单清空文本，保留格式
                for p in tc.txBody.p_lst:
                    if p.r_lst:
                        # 删除所有 run
                        for r in p.r_lst:
                            p.remove(r)
    
    if target_idx < len(table.rows):
        table.rows[target_idx]._tr.addprevious(new_tr)
    else:
        tbl.append(new_tr)
    
    return table.rows[target_idx]

def duplicate_slide(prs, index):
    """复制幻灯片，并同步形状引用的图片、标签等内部关系。"""
    source_slide = prs.slides[index]
    
    # 某些模板的当前版式包含标题占位符，python-pptx 直接复用后会让
    # 复制形状的坐标系异常。使用模板的基础版式并显式复制背景更稳定。
    new_slide = prs.slides.add_slide(prs.slide_layouts[0])
    new_slide.element.set('showMasterSp', '0')
    
    # 清除新幻灯片的默认占位符
    spTree = new_slide.shapes._spTree
    for sp in list(spTree)[2:]:
        spTree.remove(sp)
    source_spTree = source_slide.shapes._spTree
    spTree.replace(spTree[0], copy.deepcopy(source_spTree[0]))
    spTree.replace(spTree[1], copy.deepcopy(source_spTree[1]))

    source_background = source_slide.element.cSld.bg
    if source_background is None:
        source_background = source_slide.slide_layout.element.cSld.bg
    if source_background is not None:
        current_background = new_slide.element.cSld.bg
        if current_background is not None:
            new_slide.element.cSld.remove(current_background)
        new_slide.element.cSld.insert(0, copy.deepcopy(source_background))

    # 新页已经拥有自己的 slideLayout 关系；其余关系必须复制并重新分配 rId。
    # 模板表格包含 p:tags，若只复制形状，旧 rId 会误指向新页的 slideLayout，
    # 导致第二页错位，Microsoft PowerPoint 甚至会拒绝打开整个文件。
    relationship_map = {}
    for relationship in source_slide.part.rels.values():
        if relationship.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE, RT.TAGS):
            continue
        if relationship.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(
                relationship.reltype,
                relationship.target_ref,
            )
        else:
            new_rid = new_slide.part.rels.get_or_add(
                relationship.reltype,
                relationship.target_part,
            )
        relationship_map[relationship.rId] = new_rid
    
    # 复制源幻灯片的形状
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        # Office 的表格标签属于当前幻灯片的私有元数据，不能让多个页面
        # 共享同一个 tags part；复制时移除，不影响表格内容与格式。
        for custom_data in list(new_el.iter(qn('p:custDataLst'))):
            custom_data.getparent().remove(custom_data)
        for element in new_el.iter():
            for attribute_name, attribute_value in list(element.attrib.items()):
                if attribute_value in relationship_map:
                    element.set(attribute_name, relationship_map[attribute_value])
        spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_slide

def format_cell(cell, text, font_size=14, is_bold=False, color_rgb=None):
    """格式化单元格"""
    if cell is None:
        return
    tf = cell.text_frame
    tf.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        if not p.runs:
            p.add_run()
        for run in p.runs:
            run.text = text
            run.font.name = '微软雅黑'
            run.font.size = Pt(font_size)
            run.font.bold = is_bold
            run.font.color.rgb = color_rgb if color_rgb else RGBColor(0, 0, 0)

def clean_cell_merge_info(cell):
    """
    清理单元格合并信息（修复版）
    修复点：只清理垂直合并信息，保留水平合并属性。
    """
    try:
        tc = cell._tc
        if tc.tcPr is not None:
            # === 关键修复 ===
            # 只移除 vMerge，绝对不要移除 gridSpan
            v_merge = tc.tcPr.find(qn("a:vMerge"))
            if v_merge is not None:
                tc.tcPr.remove(v_merge)
    except:
        pass

def replace_text_in_slide(slide, old_txt, new_txt):
    """替换幻灯片中的文本"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if old_txt in p.text:
                    p.text = p.text.replace(old_txt, new_txt)
