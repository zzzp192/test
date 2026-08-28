import os
import tempfile
import unittest

import pandas as pd
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

import vda_processor


class VdaSummaryExtractionTests(unittest.TestCase):
    def test_excel_uses_nominal_thickness_and_not_csv_parser(self):
        source = pd.DataFrame({
            '试样编号': ['1-30m-1'],
            '实测厚度t': [None],
            '最大力Fm': [5.1747],
            '公称厚度t0': [1.6],
            '压头位移S': [3.9693],
            '角度': [30.4649],
        })

        with tempfile.TemporaryDirectory() as tmpdir:
            file_path = os.path.join(tmpdir, 'vda.xlsx')
            with pd.ExcelWriter(file_path) as writer:
                source.to_excel(writer, sheet_name='实验前后照片', index=False)
                source.to_excel(writer, sheet_name='2. VDA弯曲', index=False)

            extracted = vda_processor.standardize_vda_columns(
                vda_processor.load_vda_summary(file_path)
            )

        self.assertEqual(extracted.loc[0, 'SampleID'], '1-30m-1')
        self.assertEqual(extracted.loc[0, 'Thickness'], 1.6)
        self.assertEqual(extracted.loc[0, 'MaxForce'], 5.1747)
        self.assertEqual(extracted.loc[0, 'Displacement'], 3.9693)
        self.assertEqual(extracted.loc[0, 'Angle'], 30.4649)

    def test_two_sample_groups_keep_exact_table_boundaries(self):
        sample_ids = [
            '1-30m-1', '1-30m-2', '1-60m-1', '1-60m-2',
            '1-80m-1', '1-80m-2', '2-30m-1', '2-30m-2',
            '2-60m-1', '2-60m-2', '2-80m-1', '2-80m-2',
        ]
        source = pd.DataFrame({
            '试样编号': sample_ids,
            '实测厚度t': [None] * len(sample_ids),
            '最大力Fm': [5000.0 + idx * 100 for idx in range(len(sample_ids))],
            '公称厚度t0': [1.6] * len(sample_ids),
            '压头位移S': [3.0 + idx / 10 for idx in range(len(sample_ids))],
            '角度': [30.0 + idx for idx in range(len(sample_ids))],
        })
        project_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        template_path = os.path.join(project_dir, 'VDA弯曲角模板.pptx')

        with tempfile.TemporaryDirectory() as tmpdir:
            file_path = os.path.join(tmpdir, 'vda.xlsx')
            output_path = os.path.join(tmpdir, 'report.pptx')
            with pd.ExcelWriter(file_path) as writer:
                source.to_excel(writer, sheet_name='2. VDA弯曲', index=False)

            result = vda_processor.process_vda_report(file_path, template_path, output_path)
            presentation = Presentation(output_path)

        self.assertIn('成功生成报告', result)
        self.assertEqual(len(presentation.slides), 2)

        expected_groups = [
            ['1-30m', '1-60m', '1-80m', '2-30m'],
            ['2-60m', '2-80m'],
        ]
        for slide_index, (slide, expected) in enumerate(zip(presentation.slides, expected_groups)):
            table_shape = next(shape for shape in slide.shapes if shape.has_table)
            table = table_shape.table
            starts = [1 + index * 3 for index in range(len(expected))]
            self.assertEqual([table.cell(row, 0).text for row in starts], expected)
            self.assertEqual(table_shape.height, sum(row.height for row in table.rows))

            for row in starts:
                self.assertEqual(table.cell(row, 0)._tc.get('rowSpan'), '3')
                self.assertEqual(table.cell(row + 1, 0)._tc.get('vMerge'), '1')
                self.assertEqual(table.cell(row + 2, 0)._tc.get('vMerge'), '1')

            expected_first_force_kn = 5.0 + slide_index * 0.8
            self.assertEqual(table.cell(starts[0], 3).text, f'{expected_first_force_kn:.1f}')

            tags = table_shape.element.find('.//' + qn('p:tags'))
            if slide_index == 0:
                self.assertIsNotNone(tags)
                relationship_id = tags.get(qn('r:id'))
                self.assertEqual(slide.part.rels[relationship_id].reltype, RT.TAGS)
            else:
                self.assertIsNone(tags)


if __name__ == '__main__':
    unittest.main()
