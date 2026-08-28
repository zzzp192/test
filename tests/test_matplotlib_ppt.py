import os
import tempfile
import unittest

import pandas as pd
from PIL import Image
from openpyxl import Workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm

from curve_data import (
    load_tensile_xy_dataframe,
    load_vda_xy_dataframe,
    prepare_xy_dataframe,
    trim_tensile_tail_drop_points,
)
from matplotlib_ppt import (
    _axis_labels,
    _build_group_styles,
    _curve_group_label,
    _set_zero_origin,
    create_tensile_one_click_ppt,
    create_transparent_curve_images,
    create_vda_one_click_ppt,
    insert_curve_images,
)


class CurveDataTests(unittest.TestCase):
    def test_tensile_tail_drop_only_trims_from_drop_inside_last_fifteen_points(self):
        source = pd.DataFrame({
            "应力_1": [1500] * 10 + [1420, 1410, 1397.828747, 1394.657754, 1389.735655, 79.73718774, 0.106139352],
            "应变_1": list(range(10)) + [2.0, 2.1, 2.2776, 2.3806, 2.4504, 2.5576, 2.63],
            # This drop is outside the final-fifteen-point window and must remain untouched.
            "应力_2": [1000, 500] + [500] * 15,
            "应变_2": list(range(1, 18)),
        })

        trimmed = trim_tensile_tail_drop_points(source)

        self.assertEqual(trimmed.iloc[:, :2].dropna().shape[0], 15)
        self.assertEqual(trimmed.iloc[:, 0].dropna().iloc[-1], 1389.735655)
        self.assertEqual(trimmed.iloc[:, 1].dropna().iloc[-1], 2.4504)
        self.assertEqual(trimmed.iloc[:, 2:].dropna().shape[0], 17)

    def test_pair_swap_and_labels_match_origin_xyxy_rule(self):
        source = pd.DataFrame({
            "应力_1": [100, 200],
            "应变_1": [0.1, 0.2],
            "应力_2": [110, 210],
            "应变_2": [0.11, 0.21],
        })

        swapped = prepare_xy_dataframe(source, ["A-1", "A-2"], swap_xy=True)
        original = prepare_xy_dataframe(source, ["A-1", "A-2"], swap_xy=False)

        self.assertEqual(list(swapped.columns), ["应变_1", "A-1", "应变_2", "A-2"])
        self.assertEqual(swapped.iloc[:, 0].tolist(), [0.1, 0.2])
        self.assertEqual(swapped.iloc[:, 1].tolist(), [100, 200])
        self.assertEqual(original.iloc[:, 0].tolist(), [100, 200])
        self.assertEqual(original.iloc[:, 1].tolist(), [0.1, 0.2])

    def test_tensile_and_vda_loaders_use_shared_pair_logic(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tensile_path = os.path.join(tmpdir, "tensile.xlsx")
            with pd.ExcelWriter(tensile_path) as writer:
                pd.DataFrame({"试样编号": ["T-1", "T-2"]}).to_excel(
                    writer, sheet_name="实验报告", index=False
                )
                pd.DataFrame({
                    "应力_1": [100, 200], "应变_1": [0.1, 0.2],
                    "应力_2": [110, 210], "应变_2": [0.11, 0.21],
                }).to_excel(writer, sheet_name="原始数据", index=False)

            vda_path = os.path.join(tmpdir, "vda.xlsx")
            with pd.ExcelWriter(vda_path) as writer:
                pd.DataFrame({"试样编号": ["V-1", "V-2"]}).to_excel(
                    writer, sheet_name="2. VDA弯曲", index=False
                )
                pd.DataFrame({
                    "力_1": [1000, 2000], "位移_1": [0.1, 0.2],
                    "力_2": [2000, 3000], "位移_2": [0.2, 0.3],
                }).to_excel(writer, sheet_name="原始数据", index=False)

            tensile = load_tensile_xy_dataframe(tensile_path, swap_xy=True)
            vda = load_vda_xy_dataframe(vda_path, swap_xy=True)

        self.assertEqual(list(tensile.columns), ["应变_1", "T-1", "应变_2", "T-2"])
        self.assertEqual(list(vda.columns), ["位移_1", "V-1", "位移_2", "V-2"])
        self.assertEqual(vda["V-1"].tolist(), [1.0, 2.0])
        self.assertEqual(vda["V-2"].tolist(), [2.0, 3.0])


class MatplotlibPptTests(unittest.TestCase):
    def test_axis_labels_are_english_and_legend_labels_are_grouped(self):
        self.assertEqual(
            _axis_labels("tensile", True),
            ("Engineering strain/%", "Engineering stress/MPa"),
        )
        self.assertEqual(_axis_labels("vda", True), ("Displacement/mm", "Load/KN"))
        self.assertEqual(_curve_group_label("尾-1.4-边1-H-3"), "尾-1.4-边1-H")
        self.assertEqual(_curve_group_label("VDA-30m-A"), "VDA-30m")
        self.assertEqual(_curve_group_label("Sample A"), "Sample A")

        sample_labels = [f"Group-{group}-{sample}" for group in "ABCD" for sample in range(1, 4)]
        group_styles = _build_group_styles(sample_labels)
        self.assertEqual(list(group_styles), ["Group-A", "Group-B", "Group-C", "Group-D"])
        self.assertEqual(len(set(group_styles.values())), 4)
        self.assertEqual({line_style for _, line_style in group_styles.values()}, {"-"})

    def test_both_axes_start_at_zero(self):
        import matplotlib

        matplotlib.use("Agg")
        from matplotlib import pyplot as plt

        figure, axes = plt.subplots()
        axes.plot([2, 3], [5, 8])
        _set_zero_origin(axes)

        self.assertEqual(axes.get_xlim()[0], 0)
        self.assertEqual(axes.get_ylim()[0], 0)
        self.assertGreater(axes.get_xlim()[1], 3)
        self.assertGreater(axes.get_ylim()[1], 8)
        plt.close(figure)

    def test_png_is_transparent_and_fixed_16_by_12_cm(self):
        dataframe = pd.DataFrame({
            "X1": [0, 1, 2], "S-1": [0, 2, 3],
            "X2": [0, 1, 2], "S-2": [0, 1, 4],
        })
        with tempfile.TemporaryDirectory() as tmpdir:
            images = create_transparent_curve_images(
                dataframe,
                tmpdir,
                lines_per_graph=1,
                data_type="tensile",
                swap_xy=True,
            )
            self.assertEqual(len(images), 2)
            with Image.open(images[0]) as image:
                self.assertEqual(image.mode, "RGBA")
                self.assertLessEqual(abs(image.width - round(16 / 2.54 * 300)), 1)
                self.assertLessEqual(abs(image.height - round(12 / 2.54 * 300)), 1)
                self.assertEqual(image.getpixel((0, 0))[3], 0)

    def test_picture_is_inserted_on_slide_right_at_exact_size(self):
        dataframe = pd.DataFrame({"X": [0, 1], "S-1": [0, 2]})
        with tempfile.TemporaryDirectory() as tmpdir:
            images = create_transparent_curve_images(
                dataframe, tmpdir, 12, "tensile", True
            )
            ppt_path = os.path.join(tmpdir, "report.pptx")
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(ppt_path)

            inserted = insert_curve_images(ppt_path, images)
            result = Presentation(ppt_path)
            picture = next(
                shape for shape in result.slides[0].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )

        self.assertEqual(inserted, 1)
        self.assertEqual(picture.width, Cm(16))
        self.assertEqual(picture.height, Cm(12))
        self.assertEqual(picture.left, result.slide_width - Cm(16) - Cm(1.5))
        self.assertEqual(picture.top, (result.slide_height - Cm(12)) // 2)

    def test_tensile_one_click_generates_report_with_right_side_plot(self):
        project_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        template_path = os.path.join(project_dir, "拉伸模板.pptx")

        with tempfile.TemporaryDirectory() as tmpdir:
            source_path = os.path.join(tmpdir, "tensile.xlsx")
            output_path = os.path.join(tmpdir, "one_click.pptx")
            wb = Workbook()
            ws = wb.active
            ws.title = "实验报告"
            ws.append([
                "试样编号", "试样厚度", "规定塑性延伸强度Rp",
                "抗拉强度Rm", "最大力塑性延伸率Ag", "断裂总延伸率At",
                "断后伸长率A",
            ])
            for index in range(1, 5):
                ws.append([f"G{index}-1", 1.0, 1000, 1500, 3.0, 4.0, 5.0])
            curves = wb.create_sheet("原始数据")
            curve_headers = []
            for index in range(1, 5):
                curve_headers.extend([f"应力_{index}", f"应变_{index}"])
            curves.append(curve_headers)
            for point in range(10):
                row = []
                for index in range(1, 5):
                    row.extend([point * (100 + index), point / 100])
                curves.append(row)
            wb.save(source_path)

            message = create_tensile_one_click_ppt(
                source_path,
                template_path,
                output_path,
                lines_per_graph=12,
                swap_xy=True,
            )
            result = Presentation(output_path)
            pictures = [
                shape for shape in result.slides[0].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]

        self.assertIn("成功生成一键 PPT", message)
        self.assertEqual(len(result.slides), 1)
        self.assertEqual(len(pictures), 1)
        self.assertEqual(pictures[0].width, Cm(16))
        self.assertEqual(pictures[0].height, Cm(12))

    def test_vda_one_click_generates_report_with_right_side_plot(self):
        project_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        template_path = os.path.join(project_dir, "VDA弯曲角模板.pptx")

        sample_ids = [f"G{group}-{sample}" for group in range(1, 5) for sample in range(1, 4)]
        summary = pd.DataFrame({
            "试样编号": sample_ids,
            "公称厚度t0": [1.6] * 12,
            "最大力Fm": [5.0 + index / 10 for index in range(12)],
            "压头位移S": [3.0 + index / 10 for index in range(12)],
            "角度": [30.0 + index for index in range(12)],
        })
        raw_columns = {}
        for index, sample_id in enumerate(sample_ids, start=1):
            raw_columns[f"力_{index}"] = [point * (1 + index / 20) for point in range(10)]
            raw_columns[f"位移_{index}"] = [point / 10 for point in range(10)]

        with tempfile.TemporaryDirectory() as tmpdir:
            source_path = os.path.join(tmpdir, "vda.xlsx")
            output_path = os.path.join(tmpdir, "one_click_vda.pptx")
            with pd.ExcelWriter(source_path) as writer:
                summary.to_excel(writer, sheet_name="2. VDA弯曲", index=False)
                pd.DataFrame(raw_columns).to_excel(writer, sheet_name="原始数据", index=False)

            message = create_vda_one_click_ppt(
                source_path,
                template_path,
                output_path,
                lines_per_graph=12,
                swap_xy=True,
            )
            result = Presentation(output_path)
            pictures = [
                shape for shape in result.slides[0].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]

        self.assertIn("成功生成一键 PPT", message)
        self.assertEqual(len(result.slides), 1)
        self.assertEqual(len(pictures), 1)
        self.assertEqual(pictures[0].left, result.slide_width - Cm(16) - Cm(1.5))


if __name__ == "__main__":
    unittest.main()
