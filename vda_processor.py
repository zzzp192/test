#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
育材堂报告助手 V3.7 - VDA弯曲数据处理模块

软件名称：育材堂报告助手
版本号：V3.7
开发单位：育材堂
开发完成日期：2024年

模块功能：
    提供VDA弯曲试验数据的提取、统计计算和PPT报告生成功能。

主要功能：
    - 从Excel/CSV文件提取VDA弯曲试验数据
    - 自动识别试样编号和分组
    - 单位自动转换（N→kN）
    - 计算平均值和标准差统计
    - 动态生成PPT报告表格

数据提取字段：
    - 试样编号、公称厚度、最大力
    - 压头位移、弯曲角度

Copyright (c) 2024 育材堂. All rights reserved.
"""

import pandas as pd
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from ppt_utils import (
    to_float, delete_table_column, delete_table_row, 
    duplicate_slide, THEME_COLOR
)

REPORT_TITLE_COLOR = RGBColor(15, 78, 82)


def load_vda_summary(excel_path):
    """读取 VDA 汇总表；Excel 文件绝不按 CSV 尝试解析。"""
    extension = os.path.splitext(excel_path)[1].lower()

    if extension == '.csv':
        read_errors = []
        for encoding in ('utf-8-sig', 'utf-8', 'gb18030', 'gbk'):
            try:
                return pd.read_csv(excel_path, encoding=encoding)
            except Exception as exc:
                read_errors.append(str(exc))
        raise ValueError(f"读取 CSV 失败: {'; '.join(read_errors)}")

    xls = pd.ExcelFile(excel_path)
    try:
        preferred_sheets = [
            sheet for sheet in xls.sheet_names
            if sheet.strip() == '2. VDA弯曲'
        ]
        preferred_sheets.extend(
            sheet for sheet in xls.sheet_names
            if sheet not in preferred_sheets and 'VDA' in sheet and '原始数据' not in sheet
        )
        preferred_sheets.extend(sheet for sheet in xls.sheet_names if sheet not in preferred_sheets)

        required_headers = ('试样编号', '公称厚度', '最大力', '压头位移', '角度')
        for sheet in preferred_sheets:
            df = pd.read_excel(xls, sheet_name=sheet)
            headers = [str(column).replace('\n', '').replace(' ', '') for column in df.columns]
            if all(any(token in header for header in headers) for token in required_headers):
                return df
    finally:
        xls.close()

    raise ValueError('未找到包含试样编号、公称厚度、最大力、压头位移和角度的 VDA 汇总表。')


def standardize_vda_columns(df):
    """统一 VDA 汇总表列名；厚度始终采用公称厚度 t0。"""
    col_map = {
        "试样编号": "SampleID",
        "公称厚度t0": "Thickness",
        "最大力Fm": "MaxForce",
        "压头位移S": "Displacement",
        "角度": "Angle"
    }

    standardized = df.copy()
    for source_name, standard_name in col_map.items():
        for column in standardized.columns:
            normalized_name = str(column).replace('\n', '').replace(' ', '')
            if source_name in normalized_name:
                standardized.rename(columns={column: standard_name}, inplace=True)
                break

    return standardized

def process_vda_report(excel_path, ppt_template, output_path, force_unit='kN', include_disp=True):
    """
    处理VDA弯曲数据并生成PPT报告
    简化版：默认单位 kN，不再尝试修改表头文字
    """
    print(f"--- 开始处理 VDA 弯曲报告: {excel_path} (单位: {force_unit}) ---")
    
    # 1. 读取数据
    try:
        df = load_vda_summary(excel_path)
    except Exception as e:
        return f"读取Excel/CSV失败: {str(e)}"

    # 2. 列名映射
    df = standardize_vda_columns(df)
    
    if 'SampleID' in df.columns:
        df = df.dropna(subset=['SampleID'])
        df = df[df['SampleID'].astype(str).str.strip() != '']

    required = ["SampleID", "Thickness", "MaxForce", "Displacement", "Angle"]
    missing = [c for c in required if c not in df.columns]
    if missing:
        return f"错误: Excel中找不到这些列: {missing}，请检查表头。"

    # 原始最大力以 N 记录；报告模板使用 kN，数据行和统计值必须整体换算。
    if str(force_unit).lower() == 'kn':
        df['MaxForce'] = pd.to_numeric(df['MaxForce'], errors='coerce') / 1000.0

    # 3. 分组逻辑
    project_id = os.path.splitext(os.path.basename(excel_path))[0]
    
    def parse_group(sid):
        sid = str(sid).strip()
        if '-' in sid:
            parts = sid.rsplit('-', 1)
            return parts[0], parts[1]
        else:
            return sid, "1"

    df['GroupName'] = df['SampleID'].apply(lambda x: parse_group(x)[0])
    df['Number'] = df['SampleID'].apply(lambda x: parse_group(x)[1])

    # 4. 准备PPT页面
    prs = Presentation(ppt_template)
    unique_groups = df['GroupName'].unique()
    total_groups = len(unique_groups)
    
    if total_groups == 0:
        return "错误：未识别到任何有效的分组数据，请检查“试样编号”列。"

    groups_per_slide = 4
    num_slides_needed = (total_groups + groups_per_slide - 1) // groups_per_slide
    
    if num_slides_needed > 1:
        for _ in range(num_slides_needed - 1):
            duplicate_slide(prs, 0)
    
    group_chunks = [unique_groups[i:i + groups_per_slide] for i in range(0, total_groups, groups_per_slide)]

    # 5. 循环填充每一页
    for slide_idx, chunk_groups in enumerate(group_chunks):
        if slide_idx >= len(prs.slides): break
            
        slide = prs.slides[slide_idx]
        replace_text_in_slide(slide, "项目号", project_id)
        colorize_project_title(slide, project_id)
        
        table_shapes = [s for s in slide.shapes if s.has_table]
        if not table_shapes: continue
        main_table_shape = table_shapes[0]
        main_table = main_table_shape.table

        # 根据选项删除“压头位移”列
        if not include_disp:
            if len(main_table.columns) >= 5:
                delete_table_column(main_table, 4)
        
        # 填充数据
        process_table_chunk(main_table, chunk_groups, df, force_unit, include_disp)
        # 直接删除/新增 XML 行后，python-pptx 不会同步图形框高度。
        # 高度不一致会让 PowerPoint 把后续组显示到前一组的合并单元格中。
        main_table_shape.height = sum(row.height for row in main_table.rows)

    # 6. 保存
    try:
        prs.save(output_path)
        return f"成功生成报告！\n共 {total_groups} 组数据，{num_slides_needed} 页。\n已保存至: {output_path}"
    except Exception as e:
        return f"保存PPT失败 (请关闭已打开的同名PPT): {str(e)}"

# ================= 辅助函数 =================

def process_table_chunk(table, groups, full_df, unit, include_disp):
    HEADER_ROWS = 1
    DEFAULT_DATA_ROWS = 3
    BLOCK_SIZE = DEFAULT_DATA_ROWS + 1 
    row_offset = 0 
    
    for i in range(4):
        base_start_row = HEADER_ROWS + (i * BLOCK_SIZE)
        current_start_row = base_start_row + row_offset
        
        if i < len(groups):
            group_name = groups[i]
            group_data = full_df[full_df['GroupName'] == group_name].copy()
            
            try:
                group_data['Number'] = group_data['Number'].astype(int)
                group_data.sort_values('Number', inplace=True)
            except: pass
            
            n_data = len(group_data)
            diff = n_data - DEFAULT_DATA_ROWS
            current_stats_row_idx = current_start_row + DEFAULT_DATA_ROWS
            
            if diff > 0:
                for _ in range(diff):
                    add_table_row(table, current_stats_row_idx - 1)
                    row_offset += 1
                    current_stats_row_idx += 1
            elif diff < 0:
                rows_to_del = abs(diff)
                for _ in range(rows_to_del):
                    delete_table_row(table, current_stats_row_idx - 1)
                    current_stats_row_idx -= 1
                    row_offset -= 1
            
            fill_group_data(table, group_data, group_name, current_start_row, n_data, unit, include_disp)
            stats_idx = current_start_row + n_data
            fill_stats_row(table, group_data, stats_idx, unit, include_disp)
            
        else:
            if current_start_row >= len(table.rows): continue
            try:
                for _ in range(BLOCK_SIZE):
                    if current_start_row < len(table.rows):
                        delete_table_row(table, current_start_row)
                        row_offset -= 1
            except: pass

def fill_group_data(table, data, group_name, start_row, n_rows, unit, include_disp):
    # 模板默认每组 3 个试样加 1 行统计，第一列预先纵向合并 4 行。
    # 当实际试样数不是 3 时，删行不会自动更新 rowSpan，旧跨度会侵入下一组。
    # 先清理旧合并属性，再按“数据行 + 统计行”精确重建该组的合并区域。
    group_end_row = start_row + n_rows
    try:
        for row_idx in range(start_row, group_end_row + 1):
            cell_xml = table.cell(row_idx, 0)._tc
            cell_xml.attrib.pop('rowSpan', None)
            cell_xml.attrib.pop('vMerge', None)

        if group_end_row > start_row:
            table.cell(start_row, 0).merge(table.cell(group_end_row, 0))
    except Exception:
        pass
    
    cell_name = table.cell(start_row, 0)
    cell_name.text = str(group_name)
    format_cell(cell_name, 12)
    
    for idx, (_, row) in enumerate(data.iterrows()):
        r = start_row + idx
        format_cell_text(table, r, 1, str(row['Number']))
        format_cell_text(table, r, 2, f"{to_float(row['Thickness']):.2f}")
        
        val_f = to_float(row['MaxForce'])
        format_cell_text(table, r, 3, f"{val_f:.1f}")
        
        current_col = 4
        if include_disp:
            format_cell_text(table, r, current_col, f"{to_float(row['Displacement']):.2f}")
            current_col += 1
        
        if current_col < len(table.columns):
            format_cell_text(table, r, current_col, f"{to_float(row['Angle']):.2f}")

def fill_stats_row(table, data, r_idx, unit, include_disp):
    c_lbl = table.cell(r_idx, 1)
    c_lbl.text = "平均值±标准差"
    format_cell(c_lbl, 12, bold=True)
    custom_color = (25, 137, 141)
    
    set_stat_cell(table.cell(r_idx, 2), data['Thickness'], 1, 1.0, custom_color)
    
    f_vals = pd.to_numeric(data['MaxForce'], errors='coerce')
    set_stat_cell(table.cell(r_idx, 3), f_vals, 1, 1.0, custom_color)
    
    current_col = 4
    if include_disp:
        set_stat_cell(table.cell(r_idx, current_col), data['Displacement'], 1, 1.0, custom_color)
        current_col += 1
        
    if current_col < len(table.columns):
        set_stat_cell(table.cell(r_idx, current_col), data['Angle'], 1, 1.0, custom_color)

def format_cell_text(table, r, c, text):
    try:
        cell = table.cell(r, c)
        cell.text = text
        format_cell(cell, 12)
    except: pass

def set_stat_cell(cell, series, decimals=1, factor=1.0, color_rgb=None):
    series = pd.to_numeric(series, errors='coerce').dropna() * factor
    if len(series) == 0:
        txt = "-"
    else:
        mean = series.mean()
        std = series.std(ddof=1) if len(series) > 1 else 0.0
        txt = f"{mean:.{decimals}f}±{std:.{decimals}f}"
    
    cell.text = txt
    format_cell(cell, 12, bold=True, color=color_rgb)

def format_cell(cell, size, bold=False, color=None):
    if not cell.text_frame.paragraphs:
        cell.text_frame.text = ""
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.name = '微软雅黑'
    p.font.bold = bold
    if color:
        p.font.color.rgb = RGBColor(*color)
    p.alignment = 2

def replace_text_in_slide(slide, old_txt, new_txt):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if old_txt in p.text:
                    p.text = p.text.replace(old_txt, new_txt)

def colorize_project_title(slide, project_id):
    if not project_id:
        return
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for p in shape.text_frame.paragraphs:
            if project_id not in p.text:
                continue
            if not p.runs:
                p.add_run()
            for run in p.runs:
                run.font.color.rgb = REPORT_TITLE_COLOR

def add_table_row(table, clone_idx):
    import copy
    from pptx.oxml.ns import qn
    tr = table.rows[clone_idx]._tr
    new_tr = copy.deepcopy(tr)
    for tc in new_tr.tc_lst:
        if tc.tcPr is not None:
            for tag in ["a:vMerge", "a:gridSpan"]:
                elem = tc.tcPr.find(qn(tag))
                if elem is not None:
                    tc.tcPr.remove(elem)
    tr.addnext(new_tr)
