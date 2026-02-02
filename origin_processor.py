#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
育材堂报告助手 V3.7 - Origin绘图处理模块

软件名称：育材堂报告助手
版本号：V3.7
开发单位：育材堂
开发完成日期：2024年

模块功能：
    提供Origin绘图引擎的集成和自动化绘图功能。

主要功能：
    - Origin连接初始化和状态管理
    - 自动识别Excel数据表和试样编号
    - XY列调换和XYXY列类型设置
    - 图形创建和模板应用
    - OLE对象复制到PPT（Ctrl+J方式）
    - 剪贴板检测和重试机制
    - 拉伸曲线、VDA曲线、相变点曲线绘图
    - Origin项目文件保存

技术特点：
    - 使用originpro模块与Origin通信
    - 使用win32com操作PowerPoint
    - 使用pywin32进行Windows API调用
    - 支持自定义图片尺寸

Copyright (c) 2024 育材堂. All rights reserved.
"""

import os
import sys
import pandas as pd
import re
from pptx import Presentation
from pptx.util import Inches
import time
import win32com.client
import pythoncom
import win32gui
import win32api
import win32con
import win32clipboard
import ctypes

# 延迟导入originpro，便于错误处理
op = None

def init_origin():
    """初始化Origin连接，返回(成功, 错误信息)"""
    global op
    try:
        import originpro as _op
        op = _op
    except ImportError as e:
        return False, f"无法导入originpro模块: {e}\n请确保已安装: pip install originpro"
    except Exception as e:
        return False, f"导入originpro时出错: {e}"
    
    # 尝试连接Origin
    try:
        # 检查Origin是否运行
        if not op.oext:
            return False, "无法连接到Origin。\n\n可能原因:\n1. Origin未启动\n2. Origin版本不支持(需要Origin 2019或更高版本)\n3. Origin未执行'doc -s'命令启用Python连接\n\n解决方法:\n- 请先启动Origin 2019或更高版本\n- 在Origin中执行菜单: 连接器 > Python > 启用Python连接\n- 或在Origin命令窗口输入: doc -s"
        op.set_show(True)
        return True, None
    except Exception as e:
        err_msg = str(e)
        if "Origin" in err_msg or "COM" in err_msg:
            return False, f"连接Origin失败: {err_msg}\n\n可能原因:\n1. Origin版本过低(需要Origin 2019+)\n2. Origin未正确安装\n3. Origin未以管理员权限运行"
        return False, f"连接Origin时出错: {err_msg}"

def find_data_sheet(file_path):
    """
    自动识别数据所在的工作表
    - 包含"曲线数据"或"应变"/"应力"的为绘图表
    - 包含"试样编号"的为汇总表（提取试样编号列表）
    返回: (sheet_name, sample_ids) 
    """
    if file_path.endswith('.csv'):
        return None, None
    
    xls = pd.ExcelFile(file_path)
    data_sheet = None
    sample_ids = []
    
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        
        # 优先识别曲线数据表（名称包含"曲线"或列名包含"应变"/"应力"）
        if '曲线' in sheet_name:
            data_sheet = sheet_name
            continue
        
        # 检查列名是否包含应变/应力
        col_str = ' '.join(str(c) for c in df.columns)
        if '应变' in col_str or '应力' in col_str:
            data_sheet = sheet_name
            continue
        
        # 检查是否包含"试样编号"列（汇总表，提取试样编号列表）
        for col in df.columns:
            if '试样编号' in str(col):
                sample_ids = [str(v) for v in df[col] if pd.notna(v)]
                break
    
    return data_sheet, sample_ids

def plot_in_origin(file_path, template_path=None, lines_per_graph=1, swap_xy=False):
    """
    在 Origin 中绘图的核心函数
    """
    print(f"--- 正在尝试连接 Origin... ---")

    # --- 1. 连接 Origin ---
    try:
        if op.oext:
            op.set_show(True)
    except Exception as e:
        print(f"连接警告 (可忽略): {e}")

    # --- 2. 创建工作簿 ---
    try:
        wb = op.new_book() 
        if not wb: raise Exception("创建工作簿失败")
        wks = wb[0] 
    except Exception as e:
        return f"Origin 初始化失败: {str(e)}\n请确保 Origin 已管理员启动且执行过 doc -s"

    print(f"导入数据: {os.path.basename(file_path)}")
    
    # --- 3. 自动识别工作表 ---
    data_sheet, sample_ids = find_data_sheet(file_path)
    
    # --- 4. 导入数据 ---
    try:
        if data_sheet:
            print(f"识别到数据表: {data_sheet}")
            df = pd.read_excel(file_path, sheet_name=data_sheet)
        elif file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        
        # 如果需要交换XY列（在内存中，不修改原文件）
        if swap_xy:
            cols = list(df.columns)
            new_cols = []
            for i in range(0, len(cols) - 1, 2):
                new_cols.append(cols[i + 1])
                new_cols.append(cols[i])
            if len(cols) % 2 == 1:
                new_cols.append(cols[-1])
            df = df[new_cols]
            print(f"已交换XY列顺序")
        
        # 将试样编号设置为Y列的列名
        if sample_ids:
            cols = list(df.columns)
            for i, sid in enumerate(sample_ids):
                y_col_idx = i * 2 + 1  # Y列索引: 1, 3, 5, 7...
                if y_col_idx < len(cols):
                    cols[y_col_idx] = sid
            df.columns = cols
            print(f"已将试样编号设置为Y列名: {sample_ids[:3]}...")
        
        wks.from_df(df)
        print(f"数据导入成功，列数: {wks.cols}")
    except Exception as e:
        return f"导入数据失败: {str(e)}"

    # --- 4. 获取列数 ---
    num_cols = wks.cols
    if num_cols < 2:
        return "错误：数据列不足，至少需要 2 列 (1组 XY)"

    # --- 5. 分组绘图逻辑 (适配 XYXY) ---
    # 我们只关注 Y 列的索引：1, 3, 5, 7...
    y_cols = list(range(1, num_cols, 2))
    
    if not y_cols:
        return "警告：未找到有效的 Y 列，请检查数据列数是否为偶数。"

    chunks = [y_cols[i:i + lines_per_graph] for i in range(0, len(y_cols), lines_per_graph)]
    
    created_graphs = []

    for i, chunk in enumerate(chunks):
        template = template_path if template_path and os.path.exists(template_path) else None
        
        try:
            # 创建图形
            if template:
                graph = op.new_graph(template=template)
            else:
                graph = op.new_graph()
            layer = graph[0] 
            
            for y_idx in chunk:
                # 对应的 X 轴就在它前一列
                x_idx = y_idx - 1 
                
                # add_plot 使用整数索引
                layer.add_plot(wks, coly=y_idx, colx=x_idx, type='line')
            
            layer.rescale()
            
            fname = os.path.splitext(os.path.basename(file_path))[0]
            graph.name = f"{fname}_G{i+1}"
            created_graphs.append(graph.name)
            
        except Exception as e:
            print(f"绘图组 {i+1} 失败: {e}")

    return f"成功！\n数据模式: XYXY\n已创建 {len(created_graphs)} 张图表。\n请切换到 Origin 查看。"


def get_sample_ids_from_excel(file_path, data_type='tensile'):
    """从Excel提取试样编号列表"""
    xls = pd.ExcelFile(file_path)
    sample_ids = []
    
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        for col in df.columns:
            if '试样编号' in str(col):
                sample_ids = [str(v) for v in df[col] if pd.notna(v) and str(v).strip()]
                if sample_ids:
                    return sample_ids
    return sample_ids


def export_graph_to_image(graph, output_path):
    """导出Origin图形为EMF矢量图"""
    try:
        # 优先导出EMF格式（矢量图，可在PPT中编辑）
        emf_path = output_path.replace('.png', '.emf')
        graph.save_fig(emf_path)
        return emf_path if os.path.exists(emf_path) else None
    except:
        # 备选PNG
        try:
            graph.save_fig(output_path, width=800)
            return output_path if os.path.exists(output_path) else None
        except:
            return None


def copy_graph_to_ppt_ole(gname, prs, slide_idx, width_pt=340, height_pt=280, right_side=True):
    """使用Ctrl+J复制Origin图形为OLE对象到PPT指定页面"""
    import win32gui
    import win32api
    import win32con
    import win32clipboard
    import ctypes
    
    def get_clipboard_seq():
        """获取剪贴板序列号"""
        try:
            return ctypes.windll.user32.GetClipboardSequenceNumber()
        except:
            return 0
    
    def find_origin_window():
        result = []
        def callback(hwnd, _):
            if win32gui.IsWindowVisible(hwnd):
                title = win32gui.GetWindowText(hwnd)
                if 'Origin' in title:
                    result.append(hwnd)
            return True
        win32gui.EnumWindows(callback, None)
        return result[0] if result else None
    
    def force_origin_foreground(origin_hwnd):
        """强制Origin窗口置于前台"""
        if origin_hwnd:
            try:
                ctypes.windll.user32.AllowSetForegroundWindow(-1)
                if win32gui.IsIconic(origin_hwnd):
                    win32gui.ShowWindow(origin_hwnd, win32con.SW_RESTORE)
                win32gui.SetForegroundWindow(origin_hwnd)
                ctypes.windll.user32.SetFocus(origin_hwnd)
            except: pass
    
    def do_ctrl_j():
        """执行Ctrl+J"""
        win32api.keybd_event(win32con.VK_CONTROL, 0, 0, 0)
        time.sleep(0.1)
        win32api.keybd_event(ord('J'), 0, 0, 0)
        time.sleep(0.1)
        win32api.keybd_event(ord('J'), 0, win32con.KEYEVENTF_KEYUP, 0)
        time.sleep(0.1)
        win32api.keybd_event(win32con.VK_CONTROL, 0, win32con.KEYEVENTF_KEYUP, 0)
        time.sleep(2.5)
    
    # 清空剪贴板
    try:
        win32clipboard.OpenClipboard()
        win32clipboard.EmptyClipboard()
        win32clipboard.CloseClipboard()
    except: pass
    
    # 记录初始剪贴板序列号
    initial_seq = get_clipboard_seq()
    
    # 激活图形窗口
    op.lt_exec(f'win -a {gname};')
    time.sleep(0.5)
    
    origin_hwnd = find_origin_window()
    force_origin_foreground(origin_hwnd)
    time.sleep(0.8)
    
    # Ctrl+J 复制为OLE对象
    do_ctrl_j()
    
    # 检查剪贴板是否有变化，如果没有则重试
    new_seq = get_clipboard_seq()
    retry_count = 0
    while new_seq == initial_seq and retry_count < 3:
        print(f"[OLE] 剪贴板未变化，重试第{retry_count+1}次...")
        force_origin_foreground(origin_hwnd)
        time.sleep(0.5)
        do_ctrl_j()
        new_seq = get_clipboard_seq()
        retry_count += 1
    
    # 粘贴到PPT
    slide = prs.Slides(slide_idx)
    initial_count = slide.Shapes.Count
    try:
        slide.Shapes.Paste()
        time.sleep(0.5)
    except Exception as e:
        print(f"[OLE] 粘贴失败: {e}")
        return False
    
    # 检查是否粘贴成功，如果没有则重试
    if slide.Shapes.Count == initial_count:
        print("[OLE] 粘贴后形状数未增加，重试...")
        force_origin_foreground(origin_hwnd)
        time.sleep(0.5)
        do_ctrl_j()
        try:
            slide.Shapes.Paste()
            time.sleep(0.5)
        except: pass
    
    # 调整位置和大小
    if slide.Shapes.Count > initial_count:
        shape = slide.Shapes(slide.Shapes.Count)
        slide_width = prs.PageSetup.SlideWidth
        slide_height = prs.PageSetup.SlideHeight
        shape.Width = width_pt
        shape.Height = height_pt
        if right_side:
            shape.Left = slide_width - width_pt - 20
        else:
            shape.Left = (slide_width - width_pt) / 2
        shape.Top = (slide_height - height_pt) / 2
    
    # 保持Origin在前台
    force_origin_foreground(origin_hwnd)
    
    print(f"[OLE] 已粘贴到第{slide_idx}页: {gname}")
    return True


def create_ppt_with_origin_graphs(graph_names, output_ppt_path, folder=None):
    """导出Origin图形为PNG并创建PPT（备用方案）"""
    if folder is None:
        folder = os.path.dirname(output_ppt_path)
    
    image_paths = []
    for i, gname in enumerate(graph_names):
        img_path = os.path.join(folder, f"temp_graph_{i}.png")
        emf_path = os.path.join(folder, f"temp_graph_{i}.emf")
        try:
            graph = op.find_graph(gname)
            if graph:
                print(f"正在导出图形 {i+1}/{len(graph_names)}: {gname}")
                graph.activate()
                time.sleep(0.3)
                
                try:
                    graph.save_fig(emf_path)
                    time.sleep(0.3)
                    if os.path.exists(emf_path):
                        image_paths.append(emf_path)
                        continue
                except: pass
                
                try:
                    graph.save_fig(img_path, width=1200)
                    time.sleep(0.3)
                    if os.path.exists(img_path):
                        image_paths.append(img_path)
                except: pass
                    
        except Exception as e:
            print(f"导出图形 {gname} 失败: {e}")
    
    if image_paths:
        create_ppt_from_images(image_paths, output_ppt_path)
        for p in image_paths:
            try: os.remove(p)
            except: pass
        return output_ppt_path
    return None


def create_ppt_from_images(image_paths, output_ppt_path, origin_project_path=None):
    """将图片列表创建为PPT，每页一张图（备用方案）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 使用最后一个布局（通常是空白）
    blank_layout = prs.slide_layouts[len(prs.slide_layouts) - 1]
    
    for img_path in image_paths:
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(0.5), 
                                    width=Inches(12.333), height=Inches(6.5))
    
    prs.save(output_ppt_path)
    return output_ppt_path


def save_origin_project(output_path):
    """保存当前Origin项目"""
    try:
        op.save(output_path)
        return True
    except:
        return False


def append_origin_graphs_to_ppt(graph_names, ppt_path, folder=None):
    """将Origin图形作为EMF图片添加到已有PPT文件末尾"""
    if folder is None:
        folder = os.path.dirname(ppt_path)
    
    try:
        prs = Presentation(ppt_path)
        blank_layout = prs.slide_layouts[len(prs.slide_layouts) - 1]
        
        for i, gname in enumerate(graph_names):
            img_path = os.path.join(folder, f"temp_append_{i}.emf")
            try:
                graph = op.find_graph(gname)
                if graph:
                    graph.save_fig(img_path)
                    if os.path.exists(img_path):
                        slide = prs.slides.add_slide(blank_layout)
                        slide.shapes.add_picture(img_path, Inches(0.5), Inches(0.5), 
                                                width=Inches(12.333), height=Inches(6.5))
                        os.remove(img_path)
            except Exception as e:
                print(f"添加图形 {gname} 失败: {e}")
        
        prs.save(ppt_path)
        return ppt_path
    except Exception as e:
        print(f"添加图形到PPT失败: {e}")
        return None


def get_tensile_sample_ids(file_path):
    """从拉伸报告Excel提取试样编号列表（处理特殊格式）"""
    xls = pd.ExcelFile(file_path)
    sample_ids = []
    
    for sheet_name in xls.sheet_names:
        # 读取整个sheet，不设置header
        df_raw = pd.read_excel(xls, sheet_name=sheet_name, header=None)
        
        # 遍历每一行，找到包含"试样编号"的行
        for row_idx in range(min(10, len(df_raw))):  # 只检查前10行
            row_values = [str(v) for v in df_raw.iloc[row_idx] if pd.notna(v)]
            if any('试样编号' in v for v in row_values):
                # 找到试样编号所在的列
                for col_idx, val in enumerate(df_raw.iloc[row_idx]):
                    if pd.notna(val) and '试样编号' in str(val):
                        # 提取该列下面的所有非空值
                        sample_ids = [str(v) for v in df_raw.iloc[row_idx+1:, col_idx] 
                                     if pd.notna(v) and str(v).strip()]
                        if sample_ids:
                            print(f"从第{row_idx+1}行找到试样编号列，提取到{len(sample_ids)}个编号")
                            return sample_ids
    return sample_ids


def plot_tensile_to_ppt(file_path, template_path=None, lines_per_graph=12, swap_xy=True, append_to_ppt=None, width_cm=15.0, height_cm=12.0, copy_to_ppt=True):
    """拉伸报告Origin绘图，可选是否导出到PPT
    
    参数:
        copy_to_ppt: 是否复制图形到PPT，默认True。如果为False，只在Origin中绘图。
    """
    import win32com.client
    import pythoncom
    
    try:
        if op.oext:
            op.set_show(True)
    except: pass
    
    sample_ids = get_tensile_sample_ids(file_path)
    print(f"提取到试样编号: {sample_ids}")
    
    xls = pd.ExcelFile(file_path)
    df = None
    for sheet in xls.sheet_names:
        if '曲线' in sheet:
            df = pd.read_excel(xls, sheet_name=sheet)
            break
    if df is None:
        df = pd.read_excel(file_path)
    
    cols = list(df.columns)
    new_cols, reordered_cols = [], []
    sample_idx = 0
    for i in range(0, len(cols) - 1, 2):
        reordered_cols.extend([cols[i + 1], cols[i]])
        new_cols.append(cols[i + 1])
        new_cols.append(sample_ids[sample_idx] if sample_idx < len(sample_ids) else cols[i])
        sample_idx += 1
    if len(cols) % 2 == 1:
        reordered_cols.append(cols[-1])
        new_cols.append(cols[-1])
    
    df = df[reordered_cols]
    df.columns = new_cols
    
    wb = op.new_book()
    wks = wb[0]
    wks.from_df(df)
    
    num_cols = wks.cols
    for col_idx in range(num_cols):
        op.lt_exec(f'wks.col{col_idx + 1}.type = {4 if col_idx % 2 == 0 else 1};')
    
    y_cols = list(range(1, num_cols, 2))
    chunks = [y_cols[i:i + lines_per_graph] for i in range(0, len(y_cols), lines_per_graph)]
    
    folder = os.path.dirname(file_path)
    fname = os.path.splitext(os.path.basename(file_path))[0]
    
    # 如果不需要复制到PPT，只在Origin中绘图
    if not copy_to_ppt:
        created_graphs = []
        for i, chunk in enumerate(chunks):
            graph = op.new_graph(template=template_path) if template_path else op.new_graph()
            layer = graph[0]
            for y_idx in chunk:
                layer.add_plot(wks, coly=y_idx, colx=y_idx - 1, type='line')
            layer.rescale()
            gname = f"{fname}_T{i+1}"
            graph.name = gname
            created_graphs.append(gname)
            print(f"已完成第{i+1}/{len(chunks)}张图表")
        
        # 保存Origin项目
        opju_path = os.path.join(folder, f"拉伸曲线_{fname}.opju")
        save_origin_project(opju_path)
        
        return f"成功！已在Origin中创建 {len(created_graphs)} 张图表\nOrigin项目: {opju_path}"
    
    # 以下是复制到PPT的逻辑
    # cm转pt (1cm ≈ 28.35pt)
    width_pt = width_cm * 28.35
    height_pt = height_cm * 28.35
    
    # 初始化
    pythoncom.CoInitialize()
    op.lt_exec('doc -s;')
    time.sleep(0.3)
    
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_app.Visible = True
    
    # 如果有已存在的PPT，打开它并添加图形到每页右侧
    if append_to_ppt and os.path.exists(append_to_ppt):
        prs = ppt_app.Presentations.Open(os.path.abspath(append_to_ppt))
        ppt_path = append_to_ppt
        
        for i, chunk in enumerate(chunks):
            graph = op.new_graph(template=template_path) if template_path else op.new_graph()
            layer = graph[0]
            for y_idx in chunk:
                layer.add_plot(wks, coly=y_idx, colx=y_idx - 1, type='line')
            layer.rescale()
            gname = f"{fname}_T{i+1}"
            graph.name = gname
            
            slide_idx = i + 1
            if slide_idx <= prs.Slides.Count:
                copy_graph_to_ppt_ole(gname, prs, slide_idx, width_pt=width_pt, height_pt=height_pt, right_side=True)
            print(f"已完成第{i+1}/{len(chunks)}张图表")
    else:
        prs = ppt_app.Presentations.Add()
        ppt_path = os.path.join(folder, f"拉伸曲线_{fname}.pptx")
        
        for i, chunk in enumerate(chunks):
            graph = op.new_graph(template=template_path) if template_path else op.new_graph()
            layer = graph[0]
            for y_idx in chunk:
                layer.add_plot(wks, coly=y_idx, colx=y_idx - 1, type='line')
            layer.rescale()
            gname = f"{fname}_T{i+1}"
            graph.name = gname
            
            prs.Slides.Add(prs.Slides.Count + 1, 12)
            copy_graph_to_ppt_ole(gname, prs, prs.Slides.Count, width_pt=width_pt, height_pt=height_pt, right_side=False)
            print(f"已完成第{i+1}/{len(chunks)}张图表")
    
    opju_path = os.path.join(folder, f"拉伸曲线_{fname}.opju")
    save_origin_project(opju_path)
    
    try:
        prs.SaveAs(os.path.abspath(ppt_path))
        prs.Close()
    except: pass
    
    return f"成功！已创建 {len(chunks)} 张图表\nPPT: {ppt_path}\nOrigin项目: {opju_path}"


def plot_vda_to_ppt(file_path, template_path=None, lines_per_graph=12, swap_xy=True, width_cm=15.0, height_cm=12.0, copy_to_ppt=True):
    """VDA报告Origin绘图，可选是否导出到PPT
    
    参数:
        copy_to_ppt: 是否复制图形到PPT，默认True。如果为False，只在Origin中绘图。
    """
    import win32com.client
    import pythoncom
    
    try:
        if op.oext:
            op.set_show(True)
    except: pass
    
    sample_ids = get_sample_ids_from_excel(file_path)
    print(f"VDA提取到试样编号: {sample_ids}")
    
    xls = pd.ExcelFile(file_path)
    df = None
    for sheet in xls.sheet_names:
        if '原始数据' in sheet or 'VDA' in sheet:
            df = pd.read_excel(xls, sheet_name=sheet)
            if '力_1' in str(df.columns) or '力_' in ' '.join(str(c) for c in df.columns):
                break
    if df is None:
        df = pd.read_excel(file_path)
    
    cols = list(df.columns)
    new_cols, reordered_cols = [], []
    sample_idx = 0
    for i in range(0, len(cols) - 1, 2):
        reordered_cols.extend([cols[i + 1], cols[i]])
        new_cols.append(cols[i + 1])
        new_cols.append(sample_ids[sample_idx] if sample_idx < len(sample_ids) else cols[i])
        sample_idx += 1
    if len(cols) % 2 == 1:
        reordered_cols.append(cols[-1])
        new_cols.append(cols[-1])
    
    df = df[reordered_cols]
    df.columns = new_cols
    
    wb = op.new_book()
    wks = wb[0]
    wks.from_df(df)
    
    num_cols = wks.cols
    for col_idx in range(num_cols):
        op.lt_exec(f'wks.col{col_idx + 1}.type = {4 if col_idx % 2 == 0 else 1};')
    
    y_cols = list(range(1, num_cols, 2))
    chunks = [y_cols[i:i + lines_per_graph] for i in range(0, len(y_cols), lines_per_graph)]
    
    folder = os.path.dirname(file_path)
    fname = os.path.splitext(os.path.basename(file_path))[0]
    
    # 如果不需要复制到PPT，只在Origin中绘图
    if not copy_to_ppt:
        created_graphs = []
        for i, chunk in enumerate(chunks):
            graph = op.new_graph(template=template_path) if template_path else op.new_graph()
            layer = graph[0]
            for y_idx in chunk:
                layer.add_plot(wks, coly=y_idx, colx=y_idx - 1, type='line')
            layer.rescale()
            gname = f"{fname}_V{i+1}"
            graph.name = gname
            created_graphs.append(gname)
            print(f"已完成第{i+1}/{len(chunks)}张图表")
        
        # 保存Origin项目
        opju_path = os.path.join(folder, f"VDA曲线_{fname}.opju")
        save_origin_project(opju_path)
        
        return f"成功！已在Origin中创建 {len(created_graphs)} 张图表\nOrigin项目: {opju_path}"
    
    # 以下是复制到PPT的逻辑
    ppt_path = os.path.join(folder, f"VDA曲线_{fname}.pptx")
    
    # cm转pt (1cm ≈ 28.35pt)
    width_pt = width_cm * 28.35
    height_pt = height_cm * 28.35
    
    # 初始化PPT
    pythoncom.CoInitialize()
    op.lt_exec('doc -s;')
    time.sleep(0.3)
    
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_app.Visible = True
    prs = ppt_app.Presentations.Add()
    
    # 边绘图边导出OLE
    for i, chunk in enumerate(chunks):
        # 1. 绘图
        graph = op.new_graph(template=template_path) if template_path else op.new_graph()
        layer = graph[0]
        for y_idx in chunk:
            layer.add_plot(wks, coly=y_idx, colx=y_idx - 1, type='line')
        layer.rescale()
        gname = f"{fname}_V{i+1}"
        graph.name = gname
        
        # 2. 新建PPT页面
        prs.Slides.Add(prs.Slides.Count + 1, 12)
        
        # 3. 激活图形并Ctrl+J复制，粘贴到PPT
        copy_graph_to_ppt_ole(gname, prs, prs.Slides.Count, width_pt=width_pt, height_pt=height_pt)
        print(f"已完成第{i+1}/{len(chunks)}张图表")
    
    # 保存
    opju_path = os.path.join(folder, f"VDA曲线_{fname}.opju")
    save_origin_project(opju_path)
    
    try:
        prs.SaveAs(os.path.abspath(ppt_path))
        prs.Close()
    except: pass
    
    return f"成功！已创建 {len(chunks)} 张图表\nPPT: {ppt_path}\nOrigin项目: {opju_path}"


def find_phase_columns_from_header(file_path):
    """从CSV文件第5行表头识别温度列和长度变化列
    
    返回: (temp_col_idx, change_col_idx) 温度列索引和长度变化列索引
    """
    temp_idx = None
    change_idx = None
    
    # 读取第5行作为表头（索引为4）
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            lines = f.readlines()
            if len(lines) >= 5:
                header_line = lines[4]  # 第5行（索引4）
                # 使用分号分隔
                headers = header_line.strip().split(';')
                print(f"读取到表头: {headers}")
                
                for i, col in enumerate(headers):
                    col_str = str(col).lower().strip()
                    
                    # 识别温度列：包含 "temperature" 或 "temp" 或 "温度"
                    if temp_idx is None:
                        if 'temperature' in col_str or 'temp' in col_str or '温度' in col_str:
                            temp_idx = i
                            print(f"识别到温度列: 第{i+1}列 '{col}'")
                    
                    # 识别长度变化列：包含 "change" 或 "length" 或 "长度" 或 "变化"
                    if change_idx is None:
                        if 'change' in col_str or 'length' in col_str or '长度' in col_str or '变化' in col_str:
                            change_idx = i
                            print(f"识别到长度变化列: 第{i+1}列 '{col}'")
    except Exception as e:
        print(f"读取表头时出错: {e}")
    
    # 默认值：如果通过表头未识别到，使用默认列
    if temp_idx is None:
        temp_idx = 1  # 默认第2列
        print(f"未识别到温度列表头，使用默认: 第{temp_idx+1}列")
    if change_idx is None:
        change_idx = 3  # 默认第4列
        print(f"未识别到长度变化列表头，使用默认: 第{change_idx+1}列")
    
    return temp_idx, change_idx


def plot_phase_change(file_paths, template_path=None, width_cm=11.0, height_cm=8.8, copy_to_ppt=True):
    """相变点绘图：CSV文件，智能识别温度列和长度变化列，每个文件一张图
    
    参数:
        copy_to_ppt: 是否复制图形到PPT，默认True。如果为False，只在Origin中绘图。
    
    返回:
        如果copy_to_ppt=True: (ppt_path, opju_path, count)
        如果copy_to_ppt=False: (opju_path, count)
    """
    import win32com.client
    import pythoncom
    try:
        if op.oext:
            op.set_show(True)
    except: pass
    
    if isinstance(file_paths, str):
        file_paths = [file_paths]
    
    folder = os.path.dirname(file_paths[0]) if file_paths else "."
    count = 0
    created_graphs = []
    
    # 如果不需要复制到PPT，只在Origin中绘图
    if not copy_to_ppt:
        for i, file_path in enumerate(file_paths):
            # 先从第5行表头识别列
            temp_idx, change_idx = find_phase_columns_from_header(file_path)
            
            # 读取CSV，分号分隔，跳过前5行（第6行开始是数据，第6行是单位行也跳过）
            try:
                df = pd.read_csv(file_path, sep=';', encoding='utf-8', skiprows=6, header=None)
            except:
                try:
                    df = pd.read_csv(file_path, sep=';', encoding='gbk', skiprows=6, header=None)
                except:
                    df = pd.read_csv(file_path, skiprows=6, header=None)
            
            if df.shape[1] < 2:
                continue
            
            # 提取数据（处理欧洲格式：逗号作为小数点）
            x_data = pd.to_numeric(df.iloc[:, temp_idx].astype(str).str.replace(',', '.'), errors='coerce')
            y_data = pd.to_numeric(df.iloc[:, change_idx].astype(str).str.replace(',', '.'), errors='coerce')
            
            plot_df = pd.DataFrame({'Temperature': x_data, 'Change': y_data}).dropna()
            
            wb = op.new_book()
            wks = wb[0]
            wks.from_df(plot_df)
            
            # 绘图
            graph = op.new_graph(template=template_path) if template_path else op.new_graph()
            layer = graph[0]
            layer.add_plot(wks, coly=1, colx=0, type='line')
            layer.rescale()
            
            fname = os.path.splitext(os.path.basename(file_path))[0]
            # Origin图形命名限制24个字符，取文件名后24个字符（区分字符通常在后面）
            gname = fname[-24:] if len(fname) > 24 else fname
            graph.name = gname
            created_graphs.append(gname)
            count += 1
            print(f"已完成第{i+1}/{len(file_paths)}张图表: {fname} -> {gname}")
        
        # 保存Origin项目文件
        opju_path = os.path.join(folder, f"相变点曲线.opju")
        save_origin_project(opju_path)
        
        return opju_path, count
    
    # 以下是复制到PPT的逻辑
    # cm转pt
    width_pt = width_cm * 28.35
    height_pt = height_cm * 28.35
    
    # 初始化PPT
    pythoncom.CoInitialize()
    op.lt_exec('doc -s;')
    time.sleep(0.3)
    
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_app.Visible = True
    prs = ppt_app.Presentations.Add()
    
    ppt_path = os.path.join(folder, f"相变点曲线.pptx")
    
    # 边绘图边导出OLE
    for i, file_path in enumerate(file_paths):
        # 先从第5行表头识别列
        temp_idx, change_idx = find_phase_columns_from_header(file_path)
        
        # 读取CSV，分号分隔，跳过前6行（第5行表头+第6行单位）
        try:
            df = pd.read_csv(file_path, sep=';', encoding='utf-8', skiprows=6, header=None)
        except:
            try:
                df = pd.read_csv(file_path, sep=';', encoding='gbk', skiprows=6, header=None)
            except:
                df = pd.read_csv(file_path, skiprows=6, header=None)
        
        if df.shape[1] < 2:
            continue
        
        # 提取数据（处理欧洲格式：逗号作为小数点）
        x_data = pd.to_numeric(df.iloc[:, temp_idx].astype(str).str.replace(',', '.'), errors='coerce')
        y_data = pd.to_numeric(df.iloc[:, change_idx].astype(str).str.replace(',', '.'), errors='coerce')
        
        plot_df = pd.DataFrame({'Temperature': x_data, 'Change': y_data}).dropna()
        
        wb = op.new_book()
        wks = wb[0]
        wks.from_df(plot_df)
        
        # 绘图
        graph = op.new_graph(template=template_path) if template_path else op.new_graph()
        layer = graph[0]
        layer.add_plot(wks, coly=1, colx=0, type='line')
        layer.rescale()
        
        fname = os.path.splitext(os.path.basename(file_path))[0]
        # Origin图形命名限制24个字符，取文件名后24个字符（区分字符通常在后面）
        gname = fname[-24:] if len(fname) > 24 else fname
        graph.name = gname
        
        # 新建PPT页面并复制OLE
        prs.Slides.Add(prs.Slides.Count + 1, 12)
        copy_graph_to_ppt_ole(gname, prs, prs.Slides.Count, width_pt=width_pt, height_pt=height_pt, right_side=False)
        count += 1
        print(f"已完成第{i+1}/{len(file_paths)}张图表: {fname} -> {gname}")
    
    # 保存Origin项目文件
    opju_path = os.path.join(folder, f"相变点曲线.opju")
    save_origin_project(opju_path)
    
    try:
        prs.SaveAs(os.path.abspath(ppt_path))
        prs.Close()
    except: pass
    
    return ppt_path, opju_path, count
