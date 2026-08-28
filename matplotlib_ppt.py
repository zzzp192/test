#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Origin-free curve plotting and insertion into the report templates."""

from __future__ import annotations

import os
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

from curve_data import (
    chunk_curves,
    dataframe_to_curves,
    load_tensile_xy_dataframe,
    load_vda_xy_dataframe,
)


PLOT_WIDTH_CM = 16.0
PLOT_HEIGHT_CM = 12.0
PLOT_DPI = 300
X_AXIS_MARGIN_RATIO = 0.05
Y_AXIS_MARGIN_RATIO = 0.03
GROUP_COLOR_CYCLE = ("#000000", "#ff0000", "#00b800", "#0000ff", "#d000d0", "#00a0a0", "#ff8800", "#7030a0")
GROUP_LINE_STYLE = "-"


def _axis_labels(data_type: str, swap_xy: bool) -> tuple[str, str]:
    if data_type == "tensile":
        labels = ("Engineering strain/%", "Engineering stress/MPa")
        return labels if swap_xy else labels[::-1]
    if data_type == "vda":
        labels = ("Displacement/mm", "Load/KN")
        return labels if swap_xy else labels[::-1]
    return "X", "Y"


def _curve_group_label(sample_label: str) -> str:
    """Collapse sample suffixes such as Group-1/2/3 into one legend label."""
    label = str(sample_label).strip()
    if "-" not in label:
        return label
    group_label, sample_suffix = label.rsplit("-", 1)
    return group_label if group_label and sample_suffix else label


def _build_group_styles(sample_labels) -> dict[str, tuple[str, str]]:
    """Assign one color per group and use solid lines throughout."""
    group_styles: dict[str, tuple[str, str]] = {}
    for sample_label in sample_labels:
        group_label = _curve_group_label(sample_label)
        if group_label in group_styles:
            continue
        style_index = len(group_styles)
        group_styles[group_label] = (
            GROUP_COLOR_CYCLE[style_index % len(GROUP_COLOR_CYCLE)],
            GROUP_LINE_STYLE,
        )
    return group_styles


def _set_zero_origin(axes) -> None:
    """Start at zero and retain padded automatic upper bounds."""
    axes.margins(x=X_AXIS_MARGIN_RATIO, y=Y_AXIS_MARGIN_RATIO)
    x_upper = axes.get_xlim()[1]
    y_upper = axes.get_ylim()[1]
    axes.set_xlim(left=0, right=x_upper)
    axes.set_ylim(bottom=0, top=y_upper)


def create_transparent_curve_images(
    dataframe,
    output_dir: str,
    lines_per_graph: int,
    data_type: str,
    swap_xy: bool,
) -> list[str]:
    """Render fixed-size 16×12 cm transparent PNGs with matplotlib."""
    import matplotlib

    matplotlib.use("Agg")
    from matplotlib import pyplot as plt

    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False

    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    curve_groups = chunk_curves(dataframe_to_curves(dataframe), lines_per_graph)
    x_label, y_label = _axis_labels(data_type, swap_xy)
    image_paths: list[str] = []

    for group_index, curve_group in enumerate(curve_groups, start=1):
        figure, axes = plt.subplots(
            figsize=(PLOT_WIDTH_CM / 2.54, PLOT_HEIGHT_CM / 2.54),
            dpi=PLOT_DPI,
        )
        figure.patch.set_alpha(0)
        axes.patch.set_alpha(0)

        group_styles = _build_group_styles(curve.label for curve in curve_group)
        legend_groups: set[str] = set()
        for curve in curve_group:
            group_label = _curve_group_label(curve.label)
            color, line_style = group_styles[group_label]
            legend_label = group_label if group_label not in legend_groups else "_nolegend_"
            legend_groups.add(group_label)
            axes.plot(
                curve.x,
                curve.y,
                color=color,
                linestyle=line_style,
                linewidth=1.4,
                label=legend_label,
            )

        axes.set_xlabel(x_label, fontfamily="Times New Roman", fontsize=15)
        axes.set_ylabel(y_label, fontfamily="Times New Roman", fontsize=15)
        _set_zero_origin(axes)
        axes.minorticks_on()
        axes.tick_params(axis="both", which="major", direction="in", length=7, width=1.2, labelsize=12)
        axes.tick_params(axis="both", which="minor", direction="in", length=3.5, width=0.9)
        for tick_label in [*axes.get_xticklabels(), *axes.get_yticklabels()]:
            tick_label.set_fontfamily("Times New Roman")
        for spine in axes.spines.values():
            spine.set_color("#000000")
            spine.set_linewidth(1.4)

        if legend_groups:
            legend = axes.legend(
                loc="lower left",
                bbox_to_anchor=(0.14, 0.18),
                fontsize=11,
                frameon=False,
                handlelength=3.5,
                handletextpad=0.55,
            )
            for text in legend.get_texts():
                text.set_fontfamily("Microsoft YaHei")
        figure.subplots_adjust(left=0.14, right=0.96, bottom=0.16, top=0.965)

        image_path = output / f"curve_{group_index:03d}.png"
        figure.savefig(
            image_path,
            dpi=PLOT_DPI,
            transparent=True,
            facecolor="none",
            edgecolor="none",
        )
        plt.close(figure)
        image_paths.append(str(image_path))

    return image_paths


def insert_curve_images(
    ppt_path: str,
    image_paths: list[str],
    width_cm: float = PLOT_WIDTH_CM,
    height_cm: float = PLOT_HEIGHT_CM,
) -> int:
    """Place one fixed-size plot 1.5 cm left of the slide's right edge."""
    presentation = Presentation(ppt_path)
    if len(image_paths) > len(presentation.slides):
        raise ValueError(
            f"曲线图数量（{len(image_paths)}）多于报告页数（{len(presentation.slides)}）。"
            "请增大“每图曲线数”，确保每页最多放置一张图。"
        )

    width = Cm(width_cm)
    height = Cm(height_cm)
    left = max(0, presentation.slide_width - width - Cm(1.5))
    top = max(0, (presentation.slide_height - height) // 2)

    for slide_index, image_path in enumerate(image_paths):
        picture = presentation.slides[slide_index].shapes.add_picture(
            image_path,
            left,
            top,
            width=width,
            height=height,
        )
        picture.name = f"Matplotlib 曲线图 {slide_index + 1}"

    presentation.save(ppt_path)
    return len(image_paths)


def _create_one_click_ppt(
    data_path: str,
    template_path: str,
    output_path: str,
    lines_per_graph: int,
    swap_xy: bool,
    data_type: str,
    report_generator,
    trim_tensile_tail_drop: bool = True,
) -> str:
    if data_type == "tensile":
        dataframe = load_tensile_xy_dataframe(
            data_path,
            swap_xy=swap_xy,
            trim_tail_drop=trim_tensile_tail_drop,
        )
    else:
        dataframe = load_vda_xy_dataframe(data_path, swap_xy=swap_xy)

    output_folder = os.path.dirname(os.path.abspath(output_path))
    with tempfile.TemporaryDirectory(prefix="yucaitang_plot_", dir=output_folder) as temp_dir:
        image_paths = create_transparent_curve_images(
            dataframe,
            temp_dir,
            lines_per_graph,
            data_type,
            swap_xy,
        )
        report_message = report_generator(data_path, template_path, output_path)
        if not os.path.exists(output_path):
            raise RuntimeError(report_message or "报告生成失败")
        inserted = insert_curve_images(output_path, image_paths)

    return (
        f"成功生成一键 PPT！\n"
        f"已插入 {inserted} 张透明曲线图（12×16 cm，高×宽）。\n"
        f"保存至: {output_path}"
    )


def create_tensile_one_click_ppt(
    data_path: str,
    template_path: str,
    output_path: str,
    lines_per_graph: int = 12,
    swap_xy: bool = True,
    elongation_mode: str = "ag",
    trim_tail_drop: bool = True,
) -> str:
    import tensile_processor

    return _create_one_click_ppt(
        data_path,
        template_path,
        output_path,
        lines_per_graph,
        swap_xy,
        "tensile",
        lambda source, template, target: tensile_processor.generate_report(
            source,
            template,
            target,
            elongation_mode=elongation_mode,
        ),
        trim_tensile_tail_drop=trim_tail_drop,
    )


def create_vda_one_click_ppt(
    data_path: str,
    template_path: str,
    output_path: str,
    lines_per_graph: int = 12,
    swap_xy: bool = True,
) -> str:
    import vda_processor

    return _create_one_click_ppt(
        data_path,
        template_path,
        output_path,
        lines_per_graph,
        swap_xy,
        "vda",
        vda_processor.process_vda_report,
    )
